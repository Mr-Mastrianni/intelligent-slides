"""
Slide Deck Generator module for the Content Workflow Automation Agent.
Handles the creation, formatting, and export of slide decks.
"""
import logging
import re
import json
import os
from typing import Dict, Any, List, Optional
from datetime import datetime
import time

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    logging.warning("python-pptx not installed. PowerPoint export will not be available.")

from google_slides_exporter import GoogleSlidesExporter

# Import style templates
from style_templates import TEMPLATES

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class SlideDeckGenerator:
    """Handles the creation and formatting of slide decks."""
    
    def __init__(self):
        """Initialize the slide deck generator."""
        self.style_templates = TEMPLATES
    
    def parse_outline(self, outline_text: str) -> List[Dict[str, Any]]:
        """
        Parse the outline text into a structured format.
        
        Args:
            outline_text: The outline text to parse
            
        Returns:
            List of slide dictionaries
        """
        slides = []
        
        # Try to parse as JSON first
        try:
            # Clean up the text to handle potential code blocks in markdown
            clean_text = re.sub(r'```json\s*|\s*```', '', outline_text)
            parsed = json.loads(clean_text)
            
            # Handle different JSON formats
            if isinstance(parsed, list):
                slides = parsed
            elif isinstance(parsed, dict) and "slides" in parsed:
                slides = parsed["slides"]
            else:
                # Attempt to convert dictionary to slides
                slides = [{"title": key, "content": value} for key, value in parsed.items()]
                
        except json.JSONDecodeError:
            # If not JSON, parse as text
            logger.info("Outline is not in JSON format. Attempting to parse as text.")
            
            # Simple parsing for markdown-style headers
            current_slide = None
            for line in outline_text.split('\n'):
                line = line.strip()
                if not line:
                    continue
                
                # Check for headers (# Title)
                if line.startswith('#'):
                    if current_slide:
                        slides.append(current_slide)
                    
                    title = line.lstrip('#').strip()
                    current_slide = {
                        "title": title,
                        "content": "",
                        "points": []
                    }
                elif line.startswith('*') or line.startswith('-') or re.match(r'^\d+\.', line):
                    # Bullet points
                    point = line.lstrip('*-0123456789. ').strip()
                    if current_slide:
                        if "points" not in current_slide:
                            current_slide["points"] = []
                        current_slide["points"].append(point)
                else:
                    # Regular content
                    if current_slide:
                        if not current_slide.get("content"):
                            current_slide["content"] = line
                        else:
                            current_slide["content"] += " " + line
            
            # Add the last slide if there is one
            if current_slide:
                slides.append(current_slide)
        
        return slides
    
    def generate_slides(self, outline, style_template="default"):
        """
        Generate slides from an outline.
        
        Args:
            outline: Outline string or parsed outline
            style_template: Template to apply to slides
            
        Returns:
            List of slide dictionaries
        """
        try:
            # Parse outline if it's a string
            if isinstance(outline, str):
                parsed_outline = self.parse_outline(outline)
            else:
                parsed_outline = outline
                
            if not parsed_outline:
                logger.error("Failed to parse outline or empty outline provided")
                return []
            
            # The parse_outline method returns a list of slides
            raw_slides = parsed_outline
            formatted_slides = []
            
            # Process each slide to ensure it follows formatting guidelines
            for i, slide in enumerate(raw_slides):
                title = slide.get("title", "Slide")
                content = slide.get("content", "")
                raw_points = slide.get("points", [])
                
                # Ensure we have properly formatted points
                points = []
                for point in raw_points:
                    # Check if point already has proper format with key term and colon
                    if ":" in point and not point.startswith(":"):
                        parts = point.split(":", 1)  # Split only on the first colon
                        key_term = parts[0].strip()
                        explanation = parts[1].strip()
                        
                        # Remove any additional colons in the explanation
                        explanation = explanation.replace(":", "")
                        
                        # Capitalize each word in the key term
                        key_term = ' '.join(word.capitalize() for word in key_term.split())
                        
                        # Ensure explanation is a complete sentence
                        if explanation and not explanation.endswith(('.', '!', '?')):
                            explanation += "."
                            
                        formatted_point = f"{key_term}: {explanation}"
                        points.append(formatted_point)
                    else:
                        # Remove any colons from the content
                        clean_point = point.replace(":", "")
                        
                        # Try to format without colon by creating a proper key term
                        words = clean_point.split()
                        if len(words) >= 4:
                            # Use first 2-3 words as key term
                            key_term_length = min(3, max(2, len(words) // 4))
                            key_term = " ".join(words[:key_term_length])
                            explanation = " ".join(words[key_term_length:])
                            
                            # Capitalize each word in the key term
                            key_term = ' '.join(word.capitalize() for word in key_term.split())
                            
                            # Ensure explanation is a complete sentence
                            if explanation and not explanation.endswith(('.', '!', '?')):
                                explanation += "."
                                
                            formatted_point = f"{key_term}: {explanation}"
                            points.append(formatted_point)
                        else:
                            # For very short content, use the whole thing as an explanation
                            # with a generic key term
                            key_term = "Key Point"
                            explanation = clean_point
                            
                            # Ensure explanation is a complete sentence
                            if explanation and not explanation.endswith(('.', '!', '?')):
                                explanation += "."
                                
                            formatted_point = f"{key_term}: {explanation}"
                            points.append(formatted_point)
                
                # Make sure we have at least 5 points for each slide
                if len(points) < 5:
                    title_words = title.split()
                    clean_title = title.replace(":", "")
                    content_sentences = content.split('.')
                    
                    # Generate additional contextual points
                    additional_points_needed = 5 - len(points)
                    
                    # Extract meaningful terms from title and content for more specific key terms
                    important_terms = []
                    
                    # Add words from title that are substantive (longer than 3 chars)
                    for word in title_words:
                        if len(word) > 3 and word.lower() not in ['with', 'this', 'that', 'from', 'into', 'over', 'some', 'what']:
                            important_terms.append(word.capitalize())
                    
                    # Add any capitalized terms from content that might be meaningful
                    for sentence in content_sentences:
                        words = sentence.split()
                        for word in words:
                            clean_word = word.strip(',.;:()[]{}').capitalize()
                            if len(clean_word) > 3 and clean_word not in important_terms and clean_word[0].isupper():
                                important_terms.append(clean_word)
                    
                    # If we don't have enough important terms, add some domain-specific ones based on the title
                    if len(important_terms) < additional_points_needed:
                        # Different domain-specific terms based on the slide topic
                        if "psychedelic" in clean_title.lower() or "consciousness" in clean_title.lower():
                            domain_terms = ["Perception", "Awareness", "Cognition", "Experience", "Insight", 
                                           "Integration", "Mindfulness", "Exploration", "Neuroplasticity", "Transformation"]
                            important_terms.extend(domain_terms)
                        elif "ai" in clean_title.lower() or "artificial" in clean_title.lower():
                            domain_terms = ["Algorithm", "Data", "Learning", "Computation", "Modeling", 
                                           "Adaptation", "Inference", "Architecture", "Framework", "System"]
                            important_terms.extend(domain_terms)
                        else:
                            # General academic/business terms as fallback
                            domain_terms = ["Analysis", "Strategy", "Development", "Implementation", "Outcome", 
                                           "Research", "Innovation", "Methodology", "Framework", "Principle"]
                            important_terms.extend(domain_terms)
                    
                    # Create varied explanation templates that don't just repeat the title
                    contextual_explanations = [
                        # Specific to psychedelics and AI
                        "Creates novel connections between seemingly unrelated concepts in information processing",
                        "Enables pattern recognition beyond conventional frameworks of understanding",
                        "Facilitates enhanced learning through non-linear approaches to information",
                        "Breaks established cognitive patterns to reveal new insights and possibilities",
                        "Expands conceptual understanding beyond traditional disciplinary boundaries",
                        "Removes conventional constraints on problem-solving and ideation processes",
                        "Reveals underlying patterns typically hidden from conscious awareness",
                        "Combines analytical precision with intuitive exploration of complex phenomena",
                        "Accelerates discovery through unique approaches to knowledge integration",
                        "Offers perspectives that challenge established paradigms and methodologies",
                        "Transforms how we process information at both conscious and unconscious levels",
                        "Enhances creativity through novel recombination of existing knowledge structures"
                    ]
                    
                    # Create topic-specific explanations based on the slide content
                    topic_specific_explanations = []
                    
                    # Generate contextually relevant explanations based on the slide title
                    if "neural" in clean_title.lower() or "network" in clean_title.lower() or "pattern" in clean_title.lower():
                        topic_specific_explanations = [
                            "Neural networks exhibit remarkable patterns similar to those observed in psychedelic experiences",
                            "Provides insight into how both human brains and AI systems process pattern information",
                            "Demonstrates surprising parallels between artificial and biological processing systems",
                            "Reveals mathematical similarities between neural activation states across different systems",
                            "Shows how information processing can transcend specific hardware or wetware implementation",
                            "Highlights convergent optimization strategies in both evolved and designed systems",
                            "Suggests common underlying principles of efficient information processing",
                            "Challenges traditional distinctions between 'natural' and 'artificial' intelligence"
                        ]
                    elif "consciousness" in clean_title.lower() or "expand" in clean_title.lower() or "model" in clean_title.lower():
                        topic_specific_explanations = [
                            "Suggests consciousness may be understood as an emergent property of complex information systems",
                            "Demonstrates how both AI and psychedelics can expand beyond initial boundaries",
                            "Explores the edges of perceptual and conceptual frameworks in different systems",
                            "Offers surprising parallels between human experience expansion and AI capability growth",
                            "Proposes new frameworks for understanding the nature of consciousness itself",
                            "Challenges conventional models of mind and machine learning systems",
                            "Provides mathematical frameworks for quantifying 'expanded' states of processing",
                            "Suggests novel approaches to measuring and mapping consciousness in different systems"
                        ]
                    elif "default" in clean_title.lower() or "mode" in clean_title.lower() or "breaking" in clean_title.lower():
                        topic_specific_explanations = [
                            "Illustrates how breaking established patterns leads to novel solutions",
                            "Shows how disrupting default processing modes creates new possibilities",
                            "Demonstrates the value of controlled disorder in cognitive and computational systems",
                            "Reveals how constraints can limit innovation in both human and artificial systems",
                            "Suggests methods for intentionally disrupting established patterns to enhance creativity",
                            "Explains the neurological basis for 'breakthrough' thinking in humans and machines",
                            "Quantifies the benefits of temporary disruption of optimized processing patterns",
                            "Explores the balance between stability and innovation in complex systems"
                        ]
                    elif "therapeutic" in clean_title.lower() or "application" in clean_title.lower() or "treatment" in clean_title.lower():
                        topic_specific_explanations = [
                            "Offers promising approaches for mental health treatment using combined technologies",
                            "Demonstrates how AI can enhance the safety and efficacy of therapeutic applications",
                            "Shows potential for personalized medicine approaches to mental health",
                            "Reveals how data-driven approaches can optimize therapeutic outcomes",
                            "Suggests protocols for integrating technological and experiential treatments",
                            "Provides frameworks for responsible implementation of emerging therapies",
                            "Addresses ethical considerations in combined technological/psychedelic approaches",
                            "Explains how real-time monitoring can enhance therapeutic processes"
                        ]
                    elif "perspective" in clean_title.lower() or "novel" in clean_title.lower() or "connection" in clean_title.lower():
                        topic_specific_explanations = [
                            "Generates unexpected connections that lead to breakthrough innovations",
                            "Illustrates how cross-domain insights emerge from diverse information processing",
                            "Shows how novel perspectives arise from recombination of existing knowledge",
                            "Reveals patterns of discovery across scientific and technological advancement",
                            "Demonstrates mathematical properties of innovation in complex systems",
                            "Explains mechanisms behind creative leaps in human and artificial cognition",
                            "Provides frameworks for intentionally generating novel insights",
                            "Challenges conventional approaches to innovation and discovery"
                        ]
                    elif "ethical" in clean_title.lower() or "consideration" in clean_title.lower() or "future" in clean_title.lower():
                        topic_specific_explanations = [
                            "Raises important questions about consciousness and sentience across different systems",
                            "Explores the ethical implications of altered and artificial states of awareness",
                            "Addresses responsible development of technologies that impact fundamental cognition",
                            "Presents frameworks for evaluating benefits and risks of emergent technologies",
                            "Suggests guidelines for ethical research and application in these domains",
                            "Examines societal implications of these converging technologies",
                            "Provides balanced perspective on potential benefits and challenges",
                            "Outlines key principles for responsible progress in these fields"
                        ]
                    
                    # If we have topic-specific explanations, prioritize them
                    if topic_specific_explanations:
                        # Mix in some general ones to ensure variety
                        all_explanations = topic_specific_explanations + contextual_explanations
                    else:
                        all_explanations = contextual_explanations
                    
                    # Create a function to generate a unique seed for each slide to avoid repetition
                    slide_seed = sum(ord(c) for c in clean_title) % 100
                    
                    for j in range(additional_points_needed):
                        # Use meaningful terms we've extracted or generated
                        if important_terms and j < len(important_terms):
                            key_term = important_terms[j]
                        else:
                            # Use the term from the content if available, otherwise from title
                            content_word = f"Point {j+1}"
                            for word in content.split()[:10]:  # Look at first 10 words
                                if len(word) > 4 and word.isalpha():
                                    content_word = word.capitalize()
                                    break
                            key_term = content_word
                        
                        # Create a contextual explanation that's varied and specific
                        if j < len(content_sentences) and len(content_sentences[j].strip()) > 20:
                            # Use a sentence from the content if substantive
                            explanation = content_sentences[j].strip()
                            # Remove any colons 
                            explanation = explanation.replace(":", "")
                            # Ensure explanation is a complete sentence
                            if explanation and not explanation.endswith(('.', '!', '?')):
                                explanation += "."
                        else:
                            # Use one of our varied explanations, but with a unique index for each slide
                            # This ensures different slides get different explanations
                            explanation_index = (j + slide_seed) % len(all_explanations)
                            explanation = all_explanations[explanation_index]
                            # Ensure explanation has a period
                            if not explanation.endswith(('.', '!', '?')):
                                explanation += "."
                        
                        points.append(f"{key_term}: {explanation}")
                
                # Create the formatted slide
                formatted_slide = {
                    "title": title,
                    "content": content,
                    "points": points[:5],  # Ensure exactly 5 points
                    "type": "content" if i > 0 else "title"
                }
                formatted_slides.append(formatted_slide)
            
            return formatted_slides
            
        except Exception as e:
            logger.error(f"Error generating slides: {str(e)}")
            logger.exception(e)
            return []
    
    def format_slides(self, 
                    slides: List[Dict[str, Any]], 
                    bold_key_terms: bool = True,
                    highlight_color: Optional[str] = None) -> List[Dict[str, Any]]:
        """
        Apply formatting to slides.
        
        Args:
            slides: The slides to format
            bold_key_terms: Whether to bold key terms
            highlight_color: Optional highlight color
            
        Returns:
            List of formatted slide dictionaries
        """
        formatted_slides = []
        
        for slide in slides:
            formatted_slide = slide.copy()
            
            # Format content and points
            if bold_key_terms:
                # This is a simplified approach - in a real implementation,
                # we might use NLP to identify key terms more accurately
                
                # Bold the first few words of each bullet point
                if "points" in formatted_slide:
                    formatted_points = []
                    for point in formatted_slide["points"]:
                        words = point.split()
                        if len(words) >= 3:
                            # Bold the first word or first few words
                            bold_part = " ".join(words[:1])
                            rest_part = " ".join(words[1:])
                            formatted_point = f"**{bold_part}** {rest_part}"
                        else:
                            formatted_point = point
                        formatted_points.append(formatted_point)
                    formatted_slide["points"] = formatted_points
            
            # Apply highlight color if specified
            if highlight_color:
                formatted_slide["highlight_color"] = highlight_color
            
            formatted_slides.append(formatted_slide)
        
        logger.info(f"Formatted {len(formatted_slides)} slides")
        return formatted_slides
    
    def export_to_powerpoint(self, slides, filepath):
        """
        Export slides to a PowerPoint presentation.
        
        Args:
            slides: The slides to export
            filepath: Path to save the PowerPoint file
            
        Returns:
            Path to the exported file
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            logger.info(f"PPTX DEBUG: Starting PowerPoint creation with {len(slides)} slides")
            
            # Create a new presentation
            prs = Presentation()
            
            # Set slide dimensions to widescreen 16:9
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Get title from first slide
            title = "Presentation"
            if slides and "title" in slides[0]:
                title = slides[0]["title"]
            
            logger.info(f"PPTX DEBUG: Creating presentation with title: {title}")
            
            # Add a title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            title_shape.text = title
            subtitle_shape.text = f"Created on {datetime.now().strftime('%Y-%m-%d')}"
            
            logger.info("PPTX DEBUG: Title slide created")
            
            # Create slides (process one by one for better debugging)
            for i, slide_data in enumerate(slides):
                logger.info(f"PPTX DEBUG: Processing slide {i+1}/{len(slides)}")
                
                # Choose an appropriate layout
                layout_index = 1  # Default: Title and Content
                content_slide_layout = prs.slide_layouts[layout_index]
                
                try:
                    # Add the slide
                    slide = prs.slides.add_slide(content_slide_layout)
                    
                    # Add title if available
                    if slide.shapes.title:
                        slide.shapes.title.text = slide_data.get("title", f"Slide {i+1}")
                    
                    # Add content if placeholder exists
                    if len(slide.placeholders) > 1:
                        content = slide_data.get("content", "")
                        points = slide_data.get("points", [])
                        
                        body_shape = slide.placeholders[1]
                        tf = body_shape.text_frame
                        
                        # Add main content
                        if content:
                            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                            p.text = content
                        
                        # Add bullet points
                        for point in points:
                            p = tf.add_paragraph()
                            p.text = re.sub(r'\*\*(.*?)\*\*', r'\1', point)  # Remove markdown
                            p.level = 1
                    
                    logger.info(f"PPTX DEBUG: Slide {i+1} created successfully")
                except Exception as e:
                    logger.error(f"PPTX DEBUG: Error creating slide {i+1}: {str(e)}")
                    continue  # Continue with next slide even if this one fails
            
            # Ensure directory exists
            os.makedirs(os.path.dirname(os.path.abspath(filepath)), exist_ok=True)
            
            logger.info(f"PPTX DEBUG: Saving presentation to {filepath}")
            # Save the presentation
            prs.save(filepath)
            logger.info(f"PPTX DEBUG: Presentation saved successfully with {len(slides) + 1} slides")
            
            return filepath
        except Exception as e:
            logger.error(f"PPTX DEBUG: Fatal error in PowerPoint export: {str(e)}")
            logger.exception("Detailed PowerPoint export error:")
            raise
    
    def export_to_html(self, slides: List[Dict[str, Any]], filepath: str, style_template: str = "default") -> str:
        """
        Export slides to an HTML file that resembles modern slide presentation.
        
        Args:
            slides: List of slide dictionaries
            filepath: Path to save the HTML file
            style_template: The style template to use
            
        Returns:
            Path to the exported file
        """
        # Ensure directory exists
        os.makedirs(os.path.dirname(os.path.abspath(filepath)), exist_ok=True)
        
        # Create modern HTML representation of slides
        html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Presentation Slides</title>
    <style>
        body, html {{
            margin: 0;
            padding: 0;
            font-family: 'Segoe UI', Arial, sans-serif;
            background-color: #f5f5f5;
        }}
        .slide-container {{
            width: 100%;
            max-width: 900px;
            margin: 20px auto;
            box-shadow: 0 4px 8px rgba(0,0,0,0.1);
        }}
        .slide {{
            background-color: #444;
            color: white;
            position: relative;
            padding: 0;
            margin-bottom: 40px;
            border-radius: 4px;
            overflow: hidden;
        }}
        .slide-header {{
            padding: 30px 20px 10px;
        }}
        .slide-title {{
            font-size: 32px;
            font-weight: 600;
            margin: 0;
            padding-bottom: 10px;
        }}
        .slide-content {{
            background-color: #555;
            color: white;
            padding: 20px 30px 30px;
        }}
        .slide-main-text {{
            margin-bottom: 25px;
            font-size: 18px;
            line-height: 1.4;
        }}
        .slide-point {{
            margin: 12px 0;
            font-size: 18px;
            line-height: 1.4;
        }}
        .key-term {{
            color: #ffd700;
            font-weight: 600;
            display: inline-block;
            margin-right: 5px;
        }}
        /* Custom styles for slide types */
        .title-slide {{
            background-color: #444;
        }}
        .title-slide .slide-title {{
            font-size: 42px;
            text-align: center;
            padding: 60px 0;
        }}
        .title-slide .slide-content {{
            background-color: #444;
        }}
        .section-slide {{
            background-color: #2980b9;
        }}
        .section-slide .slide-content {{
            background-color: #3498db;
        }}
        .content-slide {{
            background-color: #444;
        }}
        .content-slide:nth-child(odd) .slide-content {{
            background-color: #505050;
        }}
        .content-slide:nth-child(even) .slide-content {{
            background-color: #555;
        }}
    </style>
</head>
<body>
    <div class="slide-container">
"""
        
        # Check if first slide should be a title slide
        if slides and len(slides) > 0:
            first_slide = slides[0]
            html_content += f"""
        <div class="slide title-slide">
            <div class="slide-header">
                <h1 class="slide-title">{first_slide['title']}</h1>
            </div>
            <div class="slide-content">
"""
            if first_slide.get('content'):
                html_content += f"""
                <div class="slide-main-text">{first_slide['content']}</div>
"""
            
            # Add bullet points for title slide if any
            for point in first_slide.get('points', []):
                # Apply highlight to the first part of each point
                formatted_point = point
                if ":" in point:
                    parts = point.split(":", 1)
                    key_term = parts[0].strip()
                    explanation = parts[1].strip()
                    formatted_point = f'<span class="key-term">{key_term}:</span> {explanation}'
                
                html_content += f"""
                <div class="slide-point">{formatted_point}</div>
"""
            
            html_content += """
            </div>
        </div>
"""
            slides = slides[1:]  # Skip the first slide in the next iteration
        
        # Generate each slide
        for i, slide in enumerate(slides):
            title = slide.get("title", "")
            content = slide.get("content", "")
            points = slide.get("points", [])
            
            # Determine slide type
            slide_class = "content-slide"
            if len(title) <= 5 or title.lower() in ["overview", "agenda", "contents", "summary", "conclusion"]:
                slide_class = "section-slide"
            
            html_content += f"""
        <div class="slide {slide_class}">
            <div class="slide-header">
                <h2 class="slide-title">{title}</h2>
            </div>
            <div class="slide-content">
"""
            # Add content if any
            if content:
                html_content += f"""
                <div class="slide-main-text">{content}</div>
"""
            
            # Add bullet points
            for point in points:
                # Apply highlight to the first part of each point
                formatted_point = point
                if ":" in point:
                    parts = point.split(":", 1)
                    key_term = parts[0].strip()
                    explanation = parts[1].strip()
                    formatted_point = f'<span class="key-term">{key_term}:</span> {explanation}'
                elif "**" in point:
                    # Convert markdown bold to HTML
                    formatted_point = point.replace("**", '<span class="key-term">', 1).replace("**", '</span>', 1)
                
                html_content += f"""
                <div class="slide-point">{formatted_point}</div>
"""
            
            html_content += """
            </div>
        </div>
"""
        
        # Close HTML
        html_content += """
    </div>
    <script>
        // Simple navigation if needed
        document.addEventListener('keydown', function(e) {
            const slides = document.querySelectorAll('.slide');
            const currentSlide = document.querySelector('.slide.active');
            let currentIndex = 0;
            
            // Initialize if no active slide
            if (!currentSlide) {
                slides[0].classList.add('active');
            } else {
                currentIndex = Array.from(slides).indexOf(currentSlide);
            }
        });
    </script>
</body>
</html>
"""
        
        # Write to file
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(html_content)
        
        logger.info(f"Exported slides to HTML: {filepath}")
        return filepath
    
    def export_to_google_slides(self, slides: List[Dict[str, Any]], filepath: str) -> str:
        """
        Export slides to Google Slides.
        
        Args:
            slides: List of slide dictionaries
            filepath: Path to save the file (used as fallback)
            
        Returns:
            URL to the Google Slides presentation or path to local HTML file as fallback
        """
        logger.info("Exporting to Google Slides")
        
        try:
            # Try to use Google Slides API integration
            google_slides_exporter = GoogleSlidesExporter(
                credentials_path=os.path.join(os.path.dirname(os.path.abspath(__file__)), 'credentials.json'),
                token_path=os.path.join(os.path.dirname(os.path.abspath(__file__)), 'token.pickle')
            )
            
            # Get title from first slide or filename
            title = None
            if slides and slides[0].get("title"):
                title = slides[0].get("title")
            else:
                # Extract title from filepath
                basename = os.path.basename(filepath)
                title = os.path.splitext(basename)[0].replace('-', ' ').title()
            
            # Export using Google Slides API
            result = google_slides_exporter.export_slides(slides, title)
            
            if result.get("status") == "success":
                logger.info(f"Successfully exported to Google Slides: {result.get('presentation_url')}")
                return result.get("presentation_url")
            else:
                logger.warning(f"Google Slides export failed: {result.get('message')}. Using HTML fallback.")
        except Exception as e:
            logger.warning(f"Error exporting to Google Slides: {str(e)}. Using HTML fallback.")
        
        # Fallback to HTML export if Google Slides export fails
        logger.info("Using HTML export with Google Slides styling as fallback")
        
        # Change file extension to .html
        if filepath.endswith('.pptx'):
            filepath = filepath.replace('.pptx', '.html')
        elif not filepath.endswith('.html'):
            filepath += '.html'
        
        return self.export_to_html(slides, filepath, style_template="google")
    
    def export_to_google_slides_html(self, slides: List[Dict[str, Any]], filepath: str) -> str:
        """
        Export to a format that resembles Google Slides (saves as HTML with Google Slides styling).
        
        Args:
            slides: List of slide dictionaries
            filepath: Path to save the file
            
        Returns:
            Path to the exported file
        """
        # Since direct Google Slides API export would require OAuth,
        # we'll create an HTML file that resembles Google Slides
        
        # Change file extension to .html
        if filepath.endswith('.pptx'):
            filepath = filepath.replace('.pptx', '.html')
        elif not filepath.endswith('.html'):
            filepath += '.html'
        
        return self.export_to_html(slides, filepath, style_template="google")
    
    def export_for_google_slides(self, slides, filepath=None):
        """
        Export slides in a format optimized for Google Slides import.
        Google Slides can import PowerPoint files (.pptx), so we'll create a simplified version
        that's more likely to import correctly.
        
        Args:
            slides: List of slide dictionaries
            filepath: Path to save the file (will default to a timestamped file if not provided)
            
        Returns:
            Dict with status and export path
        """
        logger.info("Starting export optimized for Google Slides import")
        
        # Create a default filepath if not provided
        if not filepath:
            # Get title from first slide
            title = "Presentation"
            if slides and "title" in slides[0]:
                title = slides[0]["title"]
            
            # Create a safe filename
            safe_title = re.sub(r'[^\w\-_]', '-', title.lower())
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"{safe_title}_for_google_slides_{timestamp}.pptx"
            
            # Create export directory if it doesn't exist
            export_dir = os.path.join(os.getcwd(), "exports")
            os.makedirs(export_dir, exist_ok=True)
            
            filepath = os.path.join(export_dir, filename)
        
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            logger.info(f"Creating simplified PowerPoint for Google Slides import: {filepath}")
            
            # Create a new presentation with minimal formatting
            prs = Presentation()
            
            # Set slide dimensions to widescreen 16:9 (Google Slides default)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Add a title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            # Add title and subtitle
            if slide.shapes.title:
                title_text = "Presentation"
                if slides and "title" in slides[0]:
                    title_text = slides[0]["title"]
                slide.shapes.title.text = title_text
            
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
                subtitle.text = f"Created on {datetime.now().strftime('%Y-%m-%d')}"
            
            # Add content slides with minimal formatting
            for i, slide_data in enumerate(slides):
                # Use a simple layout (Title and Content)
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                # Add title
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", f"Slide {i+1}")
                
                # Add content as bullet points
                if len(slide.placeholders) > 1:
                    content = slide_data.get("content", "")
                    points = slide_data.get("points", [])
                    
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    
                    # Add main content as introduction
                    if content:
                        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                        p.text = content
                    
                    # Add bullet points with minimal formatting
                    for point in points:
                        p = tf.add_paragraph()
                        # Remove any markdown formatting
                        clean_point = re.sub(r'\*\*(.*?)\*\*', r'\1', point)
                        p.text = clean_point
                        p.level = 1
            
            # Ensure directory exists
            os.makedirs(os.path.dirname(os.path.abspath(filepath)), exist_ok=True)
            
            # Save the presentation
            prs.save(filepath)
            logger.info(f"Successfully created PowerPoint for Google Slides import: {filepath}")
            
            return {
                "status": "success", 
                "export_path": filepath, 
                "format": "powerpoint_for_google",
                "message": "Created PowerPoint file optimized for Google Slides import"
            }
            
        except Exception as e:
            logger.error(f"Error creating PowerPoint for Google Slides: {str(e)}")
            logger.exception("Detailed error:")
            
            # Fall back to HTML if PowerPoint export fails
            html_filepath = filepath.replace('.pptx', '.html')
            try:
                export_path = self.export_to_html(slides, html_filepath)
                return {
                    "status": "partial_success", 
                    "export_path": export_path, 
                    "format": "html (fallback)",
                    "message": f"PowerPoint export failed: {str(e)}. Created HTML file instead which can be manually copied into Google Slides."
                }
            except Exception as html_error:
                return {
                    "status": "error",
                    "message": f"Failed to create PowerPoint for Google Slides: {str(e)}. HTML fallback also failed: {str(html_error)}"
                }
    
    def export_slides(self, slides, format="html", filepath=None, title=None):
        logger.info(f"EXPORT DEBUG: Starting export process with format={format}")
        logger.info(f"EXPORT DEBUG: Filepath provided: {filepath}")
        
        # Generate a default title from the first slide if not provided
        if not title and slides and "title" in slides[0]:
            title = slides[0]["title"]
        elif not title:
            title = "Slide_Deck"
        
        # Create default filename if not provided
        if not filepath:
            safe_title = re.sub(r'[^\w\-_]', '-', title.lower())
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            filename = f"{safe_title}_{timestamp}"
            
            if format == "html":
                filename = f"{filename}.html"
            elif format in ["powerpoint", "pptx"]:
                filename = f"{filename}.pptx"
            elif format == "pdf":
                filename = f"{filename}.pdf"
            elif format == "google_slides" or format == "google_slides_local":
                filename = f"{filename}_for_google.pptx"
            
            # Create export directory if it doesn't exist
            export_dir = os.path.join(os.getcwd(), "exports")
            os.makedirs(export_dir, exist_ok=True)
            
            filepath = os.path.join(export_dir, filename)
        
        logger.info(f"EXPORT DEBUG: Final filepath: {filepath}")
        
        # Continue with export based on format
        if format == "google_slides" or format == "google_slides_local":
            # Use the specialized Google Slides export
            return self.export_for_google_slides(slides, filepath)
        elif format == "powerpoint" or format == "pptx":
            try:
                logger.info("EXPORT DEBUG: Attempting PowerPoint export")
                if "Presentation" not in globals():
                    logger.info("EXPORT DEBUG: Checking for python-pptx")
                    import importlib.util
                    spec = importlib.util.find_spec("pptx")
                    if spec is None:
                        logger.warning("EXPORT DEBUG: python-pptx not found, defaulting to HTML")
                        # Fall back to HTML if pptx not available
                        return self.export_slides(slides, "html", filepath)
                
                logger.info(f"EXPORT DEBUG: PowerPoint export - slides: {len(slides)}, filepath: {filepath}")
                # Export to PowerPoint
                export_path = self.export_to_powerpoint(slides, filepath)
                logger.info(f"EXPORT DEBUG: PowerPoint export complete, path: {export_path}")
                return {"status": "success", "export_path": export_path, "format": "powerpoint"}
            except Exception as e:
                logger.error(f"EXPORT DEBUG: PowerPoint export failed with error: {str(e)}")
                logger.exception("Detailed error:")
                return {"status": "error", "message": f"PowerPoint export error: {str(e)}"}
        elif format == "html":
            try:
                export_path = self.export_to_html(slides, filepath)
                return {"status": "success", "export_path": export_path, "format": "html"}
            except Exception as e:
                logger.error(f"EXPORT DEBUG: HTML export failed with error: {str(e)}")
                return {"status": "error", "message": f"HTML export error: {str(e)}"}
        else:
            logger.error(f"EXPORT DEBUG: Unsupported format: {format}")
            return {"status": "error", "message": f"Unsupported format: {format}"}